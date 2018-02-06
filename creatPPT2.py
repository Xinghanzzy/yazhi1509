from pptx import Presentation
import sys
import random
import os


if __name__ == '__main__':
    # argv = ['test', '1', 'a',
    #         'title||two||thirteen||fourteen||D:\I.PPT 2018 @YaZhi\demopic5.png|||title||nothing||thirteen||fourteen||nothing']
    # slide_layout = argv[1]
    # path = argv[2]
    # content = argv[3]
    slide_layout = sys.argv[1]
    path = sys.argv[2]
    content = sys.argv[3]
    print(slide_layout, path, content)
    pptxname = str(random.randint(100000000, 999999999)) + ".pptx"
    # creat dir
    if os.path.isdir(path) == False:
        os.path.isdir(path)
    ospath = path + pptxname
    prs = Presentation("default_qm.pptx")
    for item_1 in content.split("|||"):
        slide = prs.slides.add_slide(prs.slide_layouts[int(slide_layout)])
        print(item_1)
        item_2 = item_1.split("||")
        # 对ppt的修改
        body_shape = slide.shapes.placeholders  # body_shape为本页ppt中所有shapes
        if item_2[0] != 'nothing':
            body_shape[0].text = item_2[0]  # Title
        if item_2[1] != 'nothing':
            body_shape[2].text = item_2[1]
        if item_2[2] != 'nothing':
            body_shape[13].text = item_2[2]
        if item_2[3] != 'nothing':
            body_shape[14].text = item_2[3]
        if item_2[4] != 'nothing':
            body_shape[1].insert_picture(item_2[4])  # Picture
    prs.save(ospath)
    print(pptxname)


'''
print(body_shape)
for shape in slide.shapes:
     print('%s' % shape.shape_type)
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
0 Title 1
1 Picture Placeholder 2
2 Text Placeholder 3
13 Text Placeholder 4
14 Text Placeholder 5

1,2
os + random.pptx

'''