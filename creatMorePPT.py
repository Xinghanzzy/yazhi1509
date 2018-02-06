from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import re

class I_ppt():
    title = description = pic = pageFrom = abstract = ""
    prs = None

    def __init__(self, title, description_top, pic, pageFrom, description_end):

        self.title = title
        self.description_top = description_top
        self.pic = pic
        self.pageFrom = pageFrom
        self.description_end = description_end
        self.prs = Presentation()

    def __init__(self):
        self.prs = Presentation()
       

    def creatPPT(self, title, description_top, pic, pageFrom, description_end):
        self.title = title
        self.description_top = description_top
        self.pic = pic
        self.pageFrom = pageFrom
        self.description_end = description_end

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Title Only

        body_shape = slide.shapes.placeholders  # body_shape为本页ppt中所有shapes
        body_shape[0].text = 'this is placeholders[0]'  # 在第一个文本框中文字框架内添加文字

        # Title
        title_shape = slide.shapes.title  # 取本页ppt的title
        title_shape.text = self.title  # 向title文本框写如文字
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(30)
            paragraph.font.bold = True

        # description_top
        top = Inches(1.23)
        left = Inches(1)
        width = Inches(9)  # 预设位置及大小
        height = Inches(1)
        textbox = slide.shapes.add_textbox(left, top, width, height)  # left，top为相对位置，width，height为文本框大小
        textbox.text = self.description_top  # 文本框中文字
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # pic
        img_path = self.pic  # 文件路径
        left, top, width, height = Inches(1.85), Inches(2.5), Inches(6), Inches(4)  # 预设位置及大小
        pic = slide.shapes.add_picture(img_path, left, top, width, height)  # 在指定位置按预设值添加图片

        # description_end
        top = Inches(6.5)
        left = Inches(1)
        width = Inches(9)  # 预设位置及大小
        height = Inches(1)
        textbox_description_end = slide.shapes.add_textbox(left, top, width, height)  # left，top为相对位置，width，height为文本框大小
        textbox_description_end.text = self.description_end  # 文本框中文字
        text_frame = textbox_description_end.text_frame
        text_frame.word_wrap = True
        for paragraph in textbox_description_end.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

        # pagefrom
        top = Inches(7)
        left = Inches(0)
        width = Inches(7)  # 预设位置及大小
        height = Inches(1)
        textbox3 = slide.shapes.add_textbox(left, top, width, height)  # left，top为相对位置，width，height为文本框大小
        textbox3.text = self.pageFrom  # 文本框中文字

    def save(self, filename):
        self.prs.save(filename)

def validateTitle(title):
    rstr = r"[\/\\\:\*\?\"\<\>\|]"  # '/ \ : * ? " < > |'
    new_title = re.sub(rstr, "_", title)  # 替换为下划线
    return new_title


list_FromXC = [{"title":"haha"},{"pageFrom":"2easd"}]
ppt = I_ppt() 
defaultPic = "C:/Users/Dante_Lab/Desktop/20180115121512.png"
for a in list_FromXC :
    title = a["title"] if a.get("title") != None and  a["title"] != "" else "Title is None"
    description_top = a["description_top"] if (a.get("description_top") != None and  a["description_top"] != "")  else "description is None"
    pageFrom = a["pageFrom"] if (a.get("pageFrom") != None and  a["pageFrom"]  != "")  else "pageFrom is None"
    description_end = a["description_end"] if (a.get("description_end") != None and  a["description_end"]  != "") else "description is None"
    pic = a["pic"] if (a.get("pic") != None and  a["pic"] != "") else defaultPic
    
    ppt.creatPPT(title, description_top, pic, pageFrom, description_end)

filename = validateTitle(title)
ppt.save("filename.pptx")




# title = "PD-L1 Status in Refractory Lymphomas"
# description_top = "'various solid tumors with marked clinical therapeutic effects due to the checkpoint blockade [anti-PD1/PD-L1 antibodies] [2], revolutionizing the treatment of solid malignancies"
# pic = "G:/QQ_Save/544130037/FileRecv/2b9d47ee42e33c0ae83d967b6688b530_r.jpg"
# pageFrom = "doi:10.1371/journal.pone.0166266.g001"
# description_end = "14:23"
# filename = 'python-pptx52.pptx'
# # a = I_ppt(title, description_top, pic, pageFrom, description_end)
# a = I_ppt() 
# a.creatPPT(title, description_top, pic, pageFrom, description_end)
# a.creatPPT(title, description_top, pic, pageFrom, description_end)
# a.creatPPT(title, description_top, pic, pageFrom, description_end)
# a.creatPPT(title, description_top, pic, pageFrom, description_end)
# a.creatPPT(title, description_top, pic, pageFrom, description_end)
# a.creatPPT(title, description_top, pic, pageFrom, description_end)
# a.save(filename)



# #3
# new_paragraph = body_shape[1].text_frame.add_paragraph()  # 在第二个shape中的文本框架中添加新段落
# new_paragraph.text = 'add_paragraph'  # 新段落中文字
# new_paragraph.font.bold = True  # 文字加粗
# new_paragraph.font.italic = True  # 文字斜体
# new_paragraph.font.size = Pt(15)  # 文字大小
# new_paragraph.font.underline = True  # 文字下划线
# new_paragraph.level = 1  # 新段落的级别
